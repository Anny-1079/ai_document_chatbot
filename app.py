import os
import streamlit as st
from dotenv import load_dotenv
from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.llms import Ollama
from langchain.retrievers.multi_query import MultiQueryRetriever
from langgraph.graph import StateGraph
from typing import TypedDict

load_dotenv()

# Constants
VECTORSTORE_DIR = "vectorstores"
os.makedirs(VECTORSTORE_DIR, exist_ok=True)

# Initialize local LLM via Ollama
llm = Ollama(model="mistral")  # Or "llama2" depending on what you pulled

# Utility function to load PPTX text
def load_pptx_text(pptx_path):
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to load or create vectorstore
def load_file_create_vectorstore(file_path, filename, ext, embeddings):
    if ext == ".pdf":
        loader = PyPDFLoader(file_path)
        pages = loader.load()
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        docs = text_splitter.split_documents(pages)
        texts = [doc.page_content for doc in docs]
    elif ext == ".pptx":
        ppt_text = load_pptx_text(file_path)
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100)
        texts = text_splitter.split_text(ppt_text)

    vectorstore_path = os.path.join(VECTORSTORE_DIR, filename)
    if os.path.exists(vectorstore_path):
        vectorstore = FAISS.load_local(vectorstore_path, embeddings, allow_dangerous_deserialization=True)
    else:
        vectorstore = FAISS.from_texts(texts, embeddings)
        vectorstore.save_local(vectorstore_path)
    return vectorstore

# Streamlit UI
st.title("📄🤖 AI Agent Document Chatbot (Ollama Local)")
uploaded_files = st.file_uploader("Upload PDFs or PPTXs", type=["pdf", "pptx"], accept_multiple_files=True)

if uploaded_files:
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    all_texts = []

    for uploaded_file in uploaded_files:
        filename, ext = os.path.splitext(uploaded_file.name)
        file_path = f"{filename}{ext}"

        with open(file_path, "wb") as f:
            f.write(uploaded_file.read())

        if ext == ".pdf":
            loader = PyPDFLoader(file_path)
            pages = loader.load()
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
            docs = text_splitter.split_documents(pages)
            texts = [doc.page_content for doc in docs]
        elif ext == ".pptx":
            ppt_text = load_pptx_text(file_path)
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100)
            texts = text_splitter.split_text(ppt_text)

        all_texts.extend(texts)

    # Create a single FAISS vectorstore from all combined texts
    combined_vectorstore = FAISS.from_texts(all_texts, embeddings)


    # MultiQuery Retriever
    multi_query_retriever = MultiQueryRetriever.from_llm(
        retriever=combined_vectorstore.as_retriever(search_kwargs={"k": 8}),
        llm=llm
    )

    # Define Graph State
    class GraphState(TypedDict):
        input: str
        result: str

    graph = StateGraph(state_schema=GraphState)

    # PDF/PPT agent node
    def agent_node(state: GraphState):
        user_input = state["input"]
        docs = multi_query_retriever.get_relevant_documents(user_input)
        context = "\n".join([doc.page_content for doc in docs])
        prompt = f"Answer based on this context:\n{context}\n\nQuestion: {user_input}"

        response = llm.invoke(prompt)
        return {"result": response, "input": user_input}

    graph.add_node("agent", agent_node)
    graph.set_entry_point("agent")

    agent_executor = graph.compile()

    user_query = st.text_input("Ask your question:")

    if user_query:
        output = agent_executor.invoke({"input": user_query})
        st.write("🤖 **Answer:**", output["result"])
